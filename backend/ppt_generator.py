from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml import parse_xml
import math
import json
import re
import os
import glob


def hex_to_rgb(hex_color):
    """Convert hex color to RGBColor"""
    if not hex_color:
        return RGBColor(255, 255, 255)
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


# Lazy-loaded dynamic themes cache
_DYNAMIC_THEMES = None

def get_dynamic_themes():
    """Get dynamic themes from database - lazy loaded"""
    global _DYNAMIC_THEMES
    if _DYNAMIC_THEMES is not None:
        return _DYNAMIC_THEMES
    
    _DYNAMIC_THEMES = {}
    try:
        from app import get_db_connection
        conn = get_db_connection()
        cursor = conn.cursor(dictionary=True)
        
        # Check if color_schemes table exists
        cursor.execute("SHOW TABLES LIKE 'tbl_color_schemes'")
        if not cursor.fetchone():
            cursor.close()
            conn.close()
            return {}
        
        # Get active color schemes
        cursor.execute("SELECT * FROM tbl_color_schemes WHERE is_active = TRUE")
        colors = cursor.fetchall()
        
        dynamic_themes = {}
        for c in colors:
            theme_name = f"dynamic_color_{c['id']}"
            dynamic_themes[theme_name] = {
                "background": hex_to_rgb(c.get('background_color', '#ffffff')),
                "title_color": hex_to_rgb(c.get('primary_color', '#007bff')),
                "body_color": hex_to_rgb(c.get('text_color', '#333333')),
                "accent_color": hex_to_rgb(c.get('accent_color', '#007bff')),
                "secondary_accent": hex_to_rgb(c.get('secondary_color', '#0056b3')),
                "gradient_start": hex_to_rgb(c.get('background_color', '#ffffff')),
                "gradient_end": hex_to_rgb(c.get('secondary_color', '#f0f0f0')),
                "shape_color": hex_to_rgb(c.get('secondary_color', '#e0e0e0')),
                "design_style": "dynamic",
                "source": "color_scheme",
                "color_scheme_id": c.get('id')
            }
        
        cursor.close()
        conn.close()
        
        _DYNAMIC_THEMES = dynamic_themes
        return dynamic_themes
    except Exception as e:
        print(f"Error loading dynamic themes: {e}")
        _DYNAMIC_THEMES = {}
        return {}


def get_template_config(theme_name):
    """Get template configuration from database by theme name (e.g., 'dynamic_1')"""
    if not theme_name.startswith('dynamic_'):
        return None
    
    try:
        # Extract template ID from theme name
        template_id = theme_name.replace('dynamic_', '')
        if not template_id.isdigit():
            return None
        
        from app import get_db_connection
        conn = get_db_connection()
        cursor = conn.cursor(dictionary=True)
        
        # Get template with its associated color scheme, font, and background
        cursor.execute("""
            SELECT t.*, 
                   cs.primary_color, cs.secondary_color, cs.accent_color, cs.background_color, cs.text_color,
                   f.font_family, f.display_name, f.font_url,
                   b.background_type, b.image_url, b.gradient_colors, b.solid_color
            FROM tbl_templates t
            LEFT JOIN tbl_color_schemes cs ON t.color_scheme_id = cs.id
            LEFT JOIN tbl_fonts f ON t.font_id = f.id
            LEFT JOIN tbl_backgrounds b ON t.background_id = b.id
            WHERE t.id = %s AND t.is_active = TRUE
        """, (template_id,))
        
        template = cursor.fetchone()
        cursor.close()
        conn.close()
        
        if not template:
            return None
        
        # Build theme config from template settings
        # Clean font_family: remove CSS fallbacks like ', sans-serif'
        raw_font = template.get('font_family', 'Arial') or 'Arial'
        clean_font = raw_font.split(',')[0].strip().strip("'\"")
        
        # Determine background image: prefer linked background image, fallback to thumbnail
        bg_type = template.get('background_type')
        bg_image = template.get('image_url')
        thumbnail = template.get('thumbnail')
        
        # If no image background linked but template has a thumbnail uploaded to theme_uploads, use it as background
        if (not bg_type or bg_type != 'image' or not bg_image) and thumbnail and '/theme_uploads/' in str(thumbnail):
            bg_type = 'image'
            bg_image = thumbnail
        
        config = {
            "background": hex_to_rgb(template.get('background_color', '#ffffff')),
            "title_color": hex_to_rgb(template.get('primary_color', '#007bff')),
            "body_color": hex_to_rgb(template.get('text_color', '#333333')),
            "accent_color": hex_to_rgb(template.get('accent_color', '#007bff')),
            "secondary_accent": hex_to_rgb(template.get('secondary_color', '#0056b3')),
            "gradient_start": hex_to_rgb(template.get('background_color', '#ffffff')),
            "gradient_end": hex_to_rgb(template.get('secondary_color', '#f0f0f0')),
            "shape_color": hex_to_rgb(template.get('secondary_color', '#e0e0e0')),
            "design_style": "template",
            "source": "template",
            "template_id": template.get('id'),
            "font_family": clean_font,
            "font_name": template.get('display_name', 'Arial'),
            "font_url": template.get('font_url', ''),
            "background_type": bg_type,
            "background_image": bg_image
        }
        
        return config
    except Exception as e:
        print(f"Error loading template config: {e}")
        return None


def clear_dynamic_themes_cache():
    """Clear the dynamic themes cache to force reload"""
    global _DYNAMIC_THEMES
    _DYNAMIC_THEMES = None


# Theme configurations with design elements
THEMES = {
    "cosmic_dark": {
        "background": RGBColor(8, 12, 28),
        "title_color": RGBColor(120, 180, 255),
        "body_color": RGBColor(220, 235, 255),
        "accent_color": RGBColor(147, 112, 219),
        "secondary_accent": RGBColor(72, 209, 204),
        "gradient_start": RGBColor(8, 12, 28),
        "gradient_end": RGBColor(18, 25, 55),
        "shape_color": RGBColor(25, 35, 70),
        "design_style": "nebula"
    },
    "gradient_blue": {
        "background": RGBColor(10, 35, 80),
        "title_color": RGBColor(255, 255, 255),
        "body_color": RGBColor(200, 230, 255),
        "accent_color": RGBColor(0, 180, 220),
        "secondary_accent": RGBColor(80, 200, 255),
        "gradient_start": RGBColor(10, 35, 80),
        "gradient_end": RGBColor(20, 60, 120),
        "shape_color": RGBColor(0, 80, 140),
        "design_style": "gradient"
    },
    "neon_purple": {
        "background": RGBColor(15, 5, 35),
        "title_color": RGBColor(220, 100, 255),
        "body_color": RGBColor(230, 180, 255),
        "accent_color": RGBColor(180, 80, 255),
        "secondary_accent": RGBColor(255, 50, 180),
        "gradient_start": RGBColor(15, 5, 35),
        "gradient_end": RGBColor(35, 15, 60),
        "shape_color": RGBColor(50, 20, 70),
        "design_style": "neon"
    },
    "minimal_light": {
        "background": RGBColor(252, 252, 255),
        "title_color": RGBColor(25, 40, 75),
        "body_color": RGBColor(55, 70, 100),
        "accent_color": RGBColor(0, 120, 215),
        "secondary_accent": RGBColor(100, 180, 255),
        "gradient_start": RGBColor(252, 252, 255),
        "gradient_end": RGBColor(245, 248, 255),
        "shape_color": RGBColor(220, 230, 245),
        "design_style": "minimal"
    },
    "corporate_pro": {
        "background": RGBColor(15, 30, 55),
        "title_color": RGBColor(255, 255, 255),
        "body_color": RGBColor(210, 225, 245),
        "accent_color": RGBColor(60, 165, 95),
        "secondary_accent": RGBColor(255, 170, 50),
        "gradient_start": RGBColor(15, 30, 55),
        "gradient_end": RGBColor(30, 50, 80),
        "shape_color": RGBColor(40, 60, 95),
        "design_style": "corporate"
    },
    "elegant_dark": {
        "background": RGBColor(12, 12, 18),
        "title_color": RGBColor(218, 180, 75),
        "body_color": RGBColor(225, 225, 225),
        "accent_color": RGBColor(185, 150, 60),
        "secondary_accent": RGBColor(218, 180, 75),
        "gradient_start": RGBColor(12, 12, 18),
        "gradient_end": RGBColor(25, 22, 30),
        "shape_color": RGBColor(35, 32, 42),
        "design_style": "elegant"
    },
    "tech_modern": {
        "background": RGBColor(10, 8, 30),
        "title_color": RGBColor(0, 255, 180),
        "body_color": RGBColor(170, 210, 255),
        "accent_color": RGBColor(80, 200, 255),
        "secondary_accent": RGBColor(0, 220, 220),
        "gradient_start": RGBColor(10, 8, 30),
        "gradient_end": RGBColor(18, 15, 45),
        "shape_color": RGBColor(25, 35, 65),
        "design_style": "tech"
    },
    "nature_fresh": {
        "background": RGBColor(12, 50, 30),
        "title_color": RGBColor(160, 245, 160),
        "body_color": RGBColor(210, 255, 200),
        "accent_color": RGBColor(90, 190, 100),
        "secondary_accent": RGBColor(180, 230, 180),
        "gradient_start": RGBColor(12, 50, 30),
        "gradient_end": RGBColor(22, 70, 45),
        "shape_color": RGBColor(35, 80, 55),
        "design_style": "nature"
    },
    "creative_colorful": {
        "background": RGBColor(25, 25, 45),
        "title_color": RGBColor(255, 130, 130),
        "body_color": RGBColor(255, 230, 170),
        "accent_color": RGBColor(255, 200, 50),
        "secondary_accent": RGBColor(130, 80, 200),
        "gradient_start": RGBColor(25, 25, 45),
        "gradient_end": RGBColor(40, 35, 60),
        "shape_color": RGBColor(55, 45, 75),
        "design_style": "creative"
    },
    "sunset_glow": {
        "background": RGBColor(35, 18, 28),
        "title_color": RGBColor(255, 200, 170),
        "body_color": RGBColor(255, 230, 215),
        "accent_color": RGBColor(255, 130, 130),
        "secondary_accent": RGBColor(255, 210, 110),
        "gradient_start": RGBColor(35, 18, 28),
        "gradient_end": RGBColor(55, 32, 42),
        "shape_color": RGBColor(70, 45, 55),
        "design_style": "sunset"
    },
    "arctic_ice": {
        "background": RGBColor(15, 40, 60),
        "title_color": RGBColor(185, 240, 225),
        "body_color": RGBColor(215, 250, 240),
        "accent_color": RGBColor(150, 225, 195),
        "secondary_accent": RGBColor(60, 110, 155),
        "gradient_start": RGBColor(15, 40, 60),
        "gradient_end": RGBColor(25, 55, 75),
        "shape_color": RGBColor(45, 70, 90),
        "design_style": "arctic"
    },
    "midnight_express": {
        "background": RGBColor(8, 8, 16),
        "title_color": RGBColor(215, 230, 255),
        "body_color": RGBColor(180, 200, 230),
        "accent_color": RGBColor(120, 150, 200),
        "secondary_accent": RGBColor(30, 45, 75),
        "gradient_start": RGBColor(8, 8, 16),
        "gradient_end": RGBColor(16, 18, 30),
        "shape_color": RGBColor(30, 35, 50),
        "design_style": "midnight"
    },
    "royal_purple": {
        "background": RGBColor(25, 10, 40),
        "title_color": RGBColor(255, 195, 230),
        "body_color": RGBColor(235, 195, 245),
        "accent_color": RGBColor(255, 70, 130),
        "secondary_accent": RGBColor(95, 55, 105),
        "gradient_start": RGBColor(25, 10, 40),
        "gradient_end": RGBColor(40, 18, 55),
        "shape_color": RGBColor(55, 28, 70),
        "design_style": "royal"
    }
}

# Font mappings — uses actual font names so the PPTX file contains exactly
# what the user selected.  If the font is not installed on the machine that
# opens the file, PowerPoint will do its own substitution (usually to a
# visually similar font), which is better than us pre-emptively replacing it.
FONTS = {
    # ── System / Office fonts ──
    "calibri": {"heading": "Calibri", "body": "Calibri"},
    "arial": {"heading": "Arial", "body": "Arial"},
    "segoe_ui": {"heading": "Segoe UI", "body": "Segoe UI"},
    "georgia": {"heading": "Georgia", "body": "Georgia"},
    "verdana": {"heading": "Verdana", "body": "Verdana"},
    "trebuchet": {"heading": "Trebuchet MS", "body": "Trebuchet MS"},
    "times": {"heading": "Times New Roman", "body": "Times New Roman"},
    "garamond": {"heading": "Garamond", "body": "Garamond"},
    "candara": {"heading": "Candara", "body": "Candara"},
    "cambria": {"heading": "Cambria", "body": "Cambria"},
    "tahoma": {"heading": "Tahoma", "body": "Tahoma"},
    "century_gothic": {"heading": "Century Gothic", "body": "Century Gothic"},
    "palatino": {"heading": "Palatino Linotype", "body": "Palatino Linotype"},
    "corbel": {"heading": "Corbel", "body": "Corbel"},
    "comic": {"heading": "Comic Sans MS", "body": "Comic Sans MS"},
    "impact": {"heading": "Impact", "body": "Impact"},
    # ── Google / web fonts (STATIC fonts only – variable fonts removed) ──
    "poppins": {"heading": "Poppins", "body": "Poppins"},
    "lato": {"heading": "Lato", "body": "Lato"},
    "fira_sans": {"heading": "Fira Sans", "body": "Fira Sans"},
    "barlow": {"heading": "Barlow", "body": "Barlow"},
    "pt_sans": {"heading": "PT Sans", "body": "PT Sans"},
    "pt_serif": {"heading": "PT Serif", "body": "PT Serif"},
    "ubuntu": {"heading": "Ubuntu", "body": "Ubuntu"},
}

# Font size configurations
FONT_SIZES = {
    "compact": {"title": 28, "subtitle": 18, "body": 12, "bullet": 11},
    "standard": {"title": 36, "subtitle": 22, "body": 16, "bullet": 14},
    "large": {"title": 44, "subtitle": 26, "body": 20, "bullet": 18},
    "extra_large": {"title": 54, "subtitle": 32, "body": 24, "bullet": 20}
}

# Custom font colors palette
FONT_COLORS = {
    "default_title": None,  # Use theme title color
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
    "title_gold": RGBColor(212, 175, 55),
    "title_blue": RGBColor(0, 150, 255),
    "title_purple": RGBColor(167, 139, 250),
    "title_green": RGBColor(76, 175, 80),
    "title_orange": RGBColor(255, 152, 0),
    "title_pink": RGBColor(255, 0, 128),
    "title_cyan": RGBColor(0, 200, 200)
}


def resolve_dynamic_color(color_key):
    """Resolve a dynamic_color_X key to an RGBColor by looking up the color scheme in the database."""
    if not color_key or not color_key.startswith('dynamic_color_'):
        return None
    try:
        color_id = color_key.replace('dynamic_color_', '')
        if not color_id.isdigit():
            return None
        from app import get_db_connection
        conn = get_db_connection()
        cursor = conn.cursor(dictionary=True)
        cursor.execute("SELECT primary_color FROM tbl_color_schemes WHERE id = %s", (color_id,))
        row = cursor.fetchone()
        cursor.close()
        conn.close()
        if row and row.get('primary_color'):
            return hex_to_rgb(row['primary_color'])
    except Exception as e:
        print(f"Warning: Could not resolve dynamic color {color_key}: {e}")
    return None


def hex_to_rgb(hex_color):
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def resolve_system_font(font_name):
    """Clean a font name for use in PowerPoint.
    Returns the actual font name (no substitution). PowerPoint will
    use the font if installed, or do its own substitution otherwise."""
    if not font_name:
        return "Calibri"
    # Clean CSS fallbacks like "Roboto, sans-serif"
    clean = font_name.split(',')[0].strip().strip("'\"")
    return clean if clean else "Calibri"


def _lookup_custom_font_from_db(font_key):
    """Look up a custom (admin-uploaded) font from the database by name.
    Returns (font_family, font_url) or (None, None) if not found."""
    try:
        from app import get_db_connection
        conn = get_db_connection()
        cursor = conn.cursor(dictionary=True)
        cursor.execute(
            "SELECT font_family, font_url FROM tbl_fonts WHERE name = %s AND is_active = TRUE LIMIT 1",
            (font_key,)
        )
        row = cursor.fetchone()
        cursor.close()
        conn.close()
        if row:
            return row.get('font_family', ''), row.get('font_url', '')
    except Exception as e:
        print(f"Warning: Could not look up custom font '{font_key}': {e}")
    return None, None


def get_font_name(font_key):
    """Get the font configuration from the font key.
    Looks up the key in the FONTS dict first; if not found, checks the
    database for admin-uploaded custom fonts; finally cleans the
    raw name and uses it directly so the PPTX contains the exact font."""
    if not font_key:
        return {"heading": "Calibri", "body": "Calibri"}
    
    # First check if it's a known key in FONTS dictionary (case-insensitive)
    font_key_lower = font_key.lower().replace(' ', '_').replace('-', '_')
    if font_key_lower in FONTS:
        return FONTS[font_key_lower]
    
    # Check database for admin-uploaded custom font
    font_family, font_url = _lookup_custom_font_from_db(font_key)
    if font_family:
        # Use the font_family from DB (the actual CSS/typeface name)
        clean = resolve_system_font(font_family)
        return {"heading": clean, "body": clean, "_font_url": font_url or ""}
    
    # Not a known key — treat as a raw font family name (e.g. from database)
    # Clean it and use as-is so PowerPoint gets the real name
    resolved = resolve_system_font(font_key)
    return {"heading": resolved, "body": resolved}


def set_gradient_background(slide, theme_config):
    """Set background on a slide - supports solid colors and background images"""
    bg_type = theme_config.get("background_type")
    bg_image = theme_config.get("background_image")

    # If the theme uses an image background, apply it as a full-slide picture
    if bg_type == 'image' and bg_image:
        try:
            # Resolve the image path on disk
            # bg_image can be a URL like "/theme_uploads/uuid_filename.png" or a full URL
            import urllib.parse
            # Handle URL-encoded characters (spaces as %20, etc.)
            decoded_image = urllib.parse.unquote(bg_image)
            image_filename = decoded_image.split('/')[-1]
            image_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'theme_uploads')
            image_path = os.path.join(image_dir, image_filename)

            # Also try with the raw (non-decoded) filename as fallback
            raw_filename = bg_image.split('/')[-1]
            raw_path = os.path.join(image_dir, raw_filename)

            # Use whichever path actually exists
            if os.path.isfile(image_path):
                final_path = image_path
            elif os.path.isfile(raw_path):
                final_path = raw_path
            else:
                # Try a glob match for partial filename match
                base_name = image_filename.split('_', 1)[-1] if '_' in image_filename else image_filename
                matches = glob.glob(os.path.join(image_dir, f'*{base_name}'))
                final_path = matches[0] if matches else None

            if final_path and os.path.isfile(final_path):
                # Use standard slide dimensions (10" x 7.5")
                slide_width = Inches(10)
                slide_height = Inches(7.5)
                # Add the image at position (0,0) covering the full slide
                pic = slide.shapes.add_picture(final_path, 0, 0, slide_width, slide_height)
                # Send the picture to the back so text/shapes render on top
                sp = pic._element
                sp.getparent().remove(sp)
                slide.shapes._spTree.insert(2, sp)
                print(f"[PPT] Background image applied: {final_path}")
                return
            else:
                print(f"Warning: Background image not found. Tried: {image_path}, {raw_path}. Falling back to solid color")
        except Exception as e:
            print(f"Warning: Could not set image background: {e}, falling back to solid color")

    # Default: solid color fill
    background = slide.background
    try:
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme_config["background"]
    except Exception as e:
        print(f"Warning: Could not set background: {e}")
        pass


def add_design_elements(slide, prs, theme_config, design_type="default"):
    """Add design elements (shapes, accents) to a slide - simplified for reliability"""
    try:
        width = prs.slide_width
        height = prs.slide_height
        accent = theme_config.get("accent_color", theme_config["title_color"])
        
        # Add accent bar at top
        try:
            top_bar = slide.shapes.add_shape(
                1,  # msoShapeRectangle
                Inches(0), Inches(0), width, Inches(0.15)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = accent
            top_bar.line.fill.background()
        except:
            pass
        
    except Exception as e:
        print(f"Warning: Could not add design elements: {e}")
        pass


def _force_run_font_xml(run, font_name):
    """Force font at XML level on a run to override any theme/default fonts.
    Sets <a:latin>, <a:cs>, and <a:ea> directly on the run properties so
    PowerPoint cannot silently substitute via theme inheritance."""
    try:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        rPr = run._r.get_or_add_rPr()

        for tag in ('latin', 'cs', 'ea'):
            elem = rPr.find('{%s}%s' % (ns, tag))
            if elem is None:
                elem = OxmlElement('a:%s' % tag)
                rPr.append(elem)
            elem.set('typeface', font_name)
            # Remove any theme-font reference that could override the typeface
            for attr in ('pitchFamily', 'charset'):
                if attr in elem.attrib:
                    del elem.attrib[attr]
    except Exception:
        pass  # python-pptx level font was already applied as fallback


def apply_text_formatting(paragraph, font_config, size_config, color_rgb=None, is_title=False, is_subtitle=False, bold=True):
    """Apply comprehensive text formatting using explicit runs for reliable font rendering"""
    # Get the font name
    if is_title:
        font_name = font_config.get("heading", "Calibri")
        font_size = size_config.get("title", 36)
    elif is_subtitle:
        font_name = font_config.get("body", "Calibri")
        font_size = size_config.get("subtitle", 22)
    else:
        font_name = font_config.get("body", "Calibri")
        font_size = size_config.get("body", 16)
    
    # Ensure text is in a Run (not just paragraph text) for reliable font application
    # If paragraph has no runs but has text, convert to a run
    if paragraph.text and len(paragraph.runs) == 0:
        text = paragraph.text
        paragraph.clear()
        run = paragraph.add_run()
        run.text = text
    elif len(paragraph.runs) == 0:
        # No text and no runs - add empty run
        run = paragraph.add_run()
    
    # Apply formatting to EVERY run (this is the reliable way in python-pptx)
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = True if is_title else bold
        if color_rgb:
            run.font.color.rgb = color_rgb
        # Force font at XML level to prevent theme/default override
        _force_run_font_xml(run, font_name)
    
    # Also set at paragraph level as fallback
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = True if is_title else bold
    if color_rgb:
        paragraph.font.color.rgb = color_rgb


def _add_corner_decorations(slide, prs, theme_config):
    """Add elegant L-shaped corner bracket decorations to a slide.
    Uses the theme's accent color (and secondary accent for inner detail)."""
    try:
        accent = theme_config.get("accent_color", theme_config["title_color"])
        secondary = theme_config.get("secondary_accent", accent)
        w = prs.slide_width   # EMU
        h = prs.slide_height  # EMU

        # Dimensions for the corner brackets
        arm_len = Inches(1.15)      # length of each arm
        thickness = Inches(0.045)   # line thickness
        margin = Inches(0.35)       # distance from slide edges
        inner_arm = Inches(0.65)    # shorter inner accent line
        inner_thick = Inches(0.025)
        inner_offset = Inches(0.14) # offset from outer bracket toward inside

        corners = [
            # (horiz_left, horiz_top, horiz_w, vert_left, vert_top, vert_h)  — TOP-LEFT
            (margin, margin, arm_len, margin, margin, arm_len),
            # TOP-RIGHT
            (w - margin - arm_len, margin, arm_len, w - margin - thickness, margin, arm_len),
            # BOTTOM-LEFT
            (margin, h - margin - thickness, arm_len, margin, h - margin - arm_len, arm_len),
            # BOTTOM-RIGHT
            (w - margin - arm_len, h - margin - thickness, arm_len,
             w - margin - thickness, h - margin - arm_len, arm_len),
        ]

        for idx, (hl, ht, hw, vl, vt, vh) in enumerate(corners):
            # Outer horizontal arm
            sh = slide.shapes.add_shape(1, hl, ht, hw, thickness)
            sh.fill.solid(); sh.fill.fore_color.rgb = accent; sh.line.fill.background()
            # Outer vertical arm
            sv = slide.shapes.add_shape(1, vl, vt, thickness, vh)
            sv.fill.solid(); sv.fill.fore_color.rgb = accent; sv.line.fill.background()

            # Inner shorter accent line (adds visual depth)
            # Offset toward the inside of the slide
            if idx == 0:    # top-left
                ih_l, ih_t = margin + inner_offset, margin + inner_offset
                iv_l, iv_t = margin + inner_offset, margin + inner_offset
            elif idx == 1:  # top-right
                ih_l, ih_t = w - margin - inner_arm - inner_offset + thickness, margin + inner_offset
                iv_l, iv_t = w - margin - thickness - inner_offset, margin + inner_offset
            elif idx == 2:  # bottom-left
                ih_l, ih_t = margin + inner_offset, h - margin - thickness - inner_offset
                iv_l, iv_t = margin + inner_offset, h - margin - inner_arm - inner_offset + thickness
            else:           # bottom-right
                ih_l, ih_t = w - margin - inner_arm - inner_offset + thickness, h - margin - thickness - inner_offset
                iv_l, iv_t = w - margin - thickness - inner_offset, h - margin - inner_arm - inner_offset + thickness

            ih = slide.shapes.add_shape(1, ih_l, ih_t, inner_arm, inner_thick)
            ih.fill.solid(); ih.fill.fore_color.rgb = secondary; ih.line.fill.background()
            # Make inner shapes semi-transparent (30% opacity)
            try:
                from lxml import etree
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                solidFill = ih._element.find('.//' + '{%s}solidFill' % ns)
                if solidFill is not None:
                    srgb = solidFill.find('{%s}srgbClr' % ns)
                    if srgb is not None:
                        alpha = OxmlElement('a:alpha')
                        alpha.set('val', '30000')  # 30%
                        srgb.append(alpha)
            except:
                pass

            iv = slide.shapes.add_shape(1, iv_l, iv_t, inner_thick, inner_arm)
            iv.fill.solid(); iv.fill.fore_color.rgb = secondary; iv.line.fill.background()
            try:
                from lxml import etree
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                solidFill = iv._element.find('.//' + '{%s}solidFill' % ns)
                if solidFill is not None:
                    srgb = solidFill.find('{%s}srgbClr' % ns)
                    if srgb is not None:
                        alpha = OxmlElement('a:alpha')
                        alpha.set('val', '30000')
                        srgb.append(alpha)
            except:
                pass

    except Exception as e:
        print(f"Warning: Could not add corner decorations: {e}")


def _add_title_accent_line(slide, prs, theme_config, y_position):
    """Add a centered decorative accent line under the title."""
    try:
        accent = theme_config.get("accent_color", theme_config["title_color"])
        line_w = Inches(2.5)
        line_h = Inches(0.035)
        left = (prs.slide_width - line_w) // 2
        sh = slide.shapes.add_shape(1, left, y_position, line_w, line_h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = accent
        sh.line.fill.background()
    except Exception as e:
        print(f"Warning: Could not add accent line: {e}")


def create_title_slide(prs, title_text, subtitle_text, theme_config, font_config, size_config):
    """Create a professional title slide with corner decorations and design elements"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    set_gradient_background(slide, theme_config)
    
    width = prs.slide_width
    height = prs.slide_height
    
    # Add elegant corner bracket decorations
    _add_corner_decorations(slide, prs, theme_config)
    
    # Add title
    try:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(p, font_config, size_config, theme_config["title_color"], is_title=True)
    except Exception as e:
        print(f"Warning: Could not add title: {e}")
    
    # Add decorative accent line beneath the title
    _add_title_accent_line(slide, prs, theme_config, Inches(4.0))
    
    # Add subtitle if provided
    if subtitle_text:
        try:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(8), Inches(1)
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle_text
            p.alignment = PP_ALIGN.CENTER
            apply_text_formatting(p, font_config, size_config, theme_config["body_color"], is_subtitle=True)
        except Exception as e:
            print(f"Warning: Could not add subtitle: {e}")
    
    return slide


def create_thankyou_slide(prs, theme_config, font_config, size_config, thank_you_text="Thank You!"):
    """Create a professional Thank You ending slide with corner decorations"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    set_gradient_background(slide, theme_config)
    
    # Add elegant corner bracket decorations (same as title slide)
    _add_corner_decorations(slide, prs, theme_config)
    
    # Add Thank You text - centered and large
    try:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.3), Inches(9), Inches(1.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = thank_you_text
        p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(p, font_config, size_config, theme_config["title_color"], is_title=True)
    except Exception as e:
        print(f"Warning: Could not add thank you title: {e}")
    
    # Add decorative accent line beneath the text
    _add_title_accent_line(slide, prs, theme_config, Inches(4.0))
    
    # Add a small tagline beneath
    try:
        tagline_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(4.3), Inches(7), Inches(0.8)
        )
        tf = tagline_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "We appreciate your attention"
        p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(p, font_config, size_config, theme_config["body_color"], is_subtitle=True)
    except Exception as e:
        print(f"Warning: Could not add thank you tagline: {e}")
    
    return slide


def create_content_slide(prs, title_text, bullets, theme_config, font_config, size_config, layout_type="bullets"):
    """Create a professional content slide with design elements"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    set_gradient_background(slide, theme_config)
    
    # Add design elements
    add_design_elements(slide, prs, theme_config)
    
    # --- DYNAMIC TITLE SIZING ---
    title_text_clean = title_text.strip()
    # Truncate extremely long titles
    if len(title_text_clean) > 100:
        title_text_clean = title_text_clean[:97] + "..."
    
    # Realistic chars per line at 36pt font on 9" wide box
    title_font_size = size_config.get("title", 36)
    if title_font_size >= 44:
        chars_per_line = 22
    elif title_font_size >= 36:
        chars_per_line = 28
    elif title_font_size >= 28:
        chars_per_line = 35
    else:
        chars_per_line = 42
    
    title_line_count = max(1, -(-len(title_text_clean) // chars_per_line))  # ceiling division
    
    # Title box: fixed area at the top
    title_box_height = min(Inches(1.8), max(Inches(0.7), Inches(0.55) * title_line_count))
    title_top = Inches(0.25)
    
    # Add title
    try:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), title_top, Inches(9), title_box_height
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.text = title_text_clean
        apply_text_formatting(p, font_config, size_config, theme_config["title_color"], is_title=True)
    except Exception as e:
        print(f"Warning: Could not add slide title: {e}")
    
    # Limit bullets to prevent overflow
    max_bullets = 6
    if len(bullets) > max_bullets:
        bullets = bullets[:max_bullets]
    
    # Content area: spans the middle-to-bottom zone of the slide
    content_top_start = title_top + title_box_height + Inches(0.2)
    slide_bottom_margin = Inches(0.5)
    available_height = prs.slide_height - content_top_start - slide_bottom_margin
    
    # Estimate how much vertical space the bullets actually need
    body_font_size = size_config.get("body", 16)
    num_total = len(bullets)
    
    avg_bullet_len = sum(len(b) for b in bullets) / max(num_total, 1)
    if body_font_size >= 20:
        chars_per_body_line = 55
    elif body_font_size >= 16:
        chars_per_body_line = 65
    else:
        chars_per_body_line = 80
    
    avg_lines_per_bullet = max(1, -(-int(avg_bullet_len) // chars_per_body_line))
    line_height_pts = body_font_size * 1.3
    spacing_pts = 14 if num_total <= 3 else (11 if num_total <= 4 else (8 if num_total <= 5 else 6))
    per_bullet_height = (avg_lines_per_bullet * line_height_pts + spacing_pts)
    estimated_content_pts = per_bullet_height * num_total
    estimated_content_inches = estimated_content_pts / 72.0
    
    # Center the content box vertically in the available space
    available_height_inches = available_height / 914400  # EMU to inches
    if estimated_content_inches < available_height_inches * 0.85:
        # Content is smaller than available space - center it vertically
        vertical_pad = (available_height_inches - estimated_content_inches) / 2
        content_top = content_top_start + Inches(vertical_pad)
        content_height = Inches(estimated_content_inches + 0.5)  # buffer
    else:
        # Content fills most of the space - use full area
        content_top = content_top_start
        content_height = available_height
    
    # Add content based on layout type
    if layout_type == "bullets":
        try:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), content_top, Inches(9), content_height
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Clean the bullet text - keep full text, no hard truncation
                clean_text = bullet.replace("\u2022", "").replace("\u2013", "-").replace("\u2014", "-").strip()
                # Only truncate extremely long bullets (>400 chars)
                if len(clean_text) > 400:
                    cut = clean_text[:397].rfind(' ')
                    if cut > 300:
                        clean_text = clean_text[:cut] + "..."
                    else:
                        clean_text = clean_text[:397] + "..."
                p.text = "\u2022 " + clean_text
                p.level = 0
                # Dynamic spacing based on bullet count
                if num_total <= 3:
                    p.space_before = Pt(16)
                    p.space_after = Pt(10)
                elif num_total <= 4:
                    p.space_before = Pt(12)
                    p.space_after = Pt(8)
                elif num_total <= 5:
                    p.space_before = Pt(8)
                    p.space_after = Pt(5)
                else:
                    p.space_before = Pt(6)
                    p.space_after = Pt(3)
                
                apply_text_formatting(p, font_config, size_config, theme_config["body_color"], is_title=False)
        except Exception as e:
            print(f"Warning: Could not add bullet content: {e}")
    else:
        # Two column layout
        # Split bullets into two columns first
        mid = len(bullets) // 2
        left_bullets = bullets[:mid]
        right_bullets = bullets[mid:]
        
        # Calculate optimal content height based on bullet count
        num_bullets = len(bullets) if bullets else 0
        col_height = min(Inches(5), Inches(1.5) + (max(len(left_bullets), len(right_bullets)) * Inches(0.35)))
        content_top = Inches(1.3)
        
        try:
            left_box = slide.shapes.add_textbox(
                Inches(0.3), content_top, Inches(4.3), col_height
            )
            right_box = slide.shapes.add_textbox(
                Inches(5.2), content_top, Inches(4.3), col_height
            )
            
            # Left column
            tf_left = left_box.text_frame
            tf_left.word_wrap = True
            try:
                tf_left.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except:
                pass
            for i, bullet in enumerate(left_bullets):
                if i == 0:
                    p = tf_left.paragraphs[0]
                else:
                    p = tf_left.add_paragraph()
                clean_text = bullet.replace("•", "").replace("–", "-").replace("—", "-").strip()
                # Truncate long bullets for readability
                if len(clean_text) > 80:
                    clean_text = clean_text[:77] + "..."
                p.text = clean_text
                p.level = 0
                p.space_before = Pt(8)
                p.space_after = Pt(4)
                apply_text_formatting(p, font_config, size_config, theme_config["body_color"], is_title=False)
            
            # Right column
            tf_right = right_box.text_frame
            tf_right.word_wrap = True
            try:
                tf_right.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except:
                pass
            for i, bullet in enumerate(right_bullets):
                if i == 0:
                    p = tf_right.paragraphs[0]
                else:
                    p = tf_right.add_paragraph()
                clean_text = bullet.replace("•", "").replace("–", "-").replace("—", "-").strip()
                # Truncate long bullets for readability
                if len(clean_text) > 80:
                    clean_text = clean_text[:77] + "..."
                p.text = clean_text
                p.level = 0
                p.space_before = Pt(8)
                p.space_after = Pt(4)
                apply_text_formatting(p, font_config, size_config, theme_config["body_color"], is_title=False)
        except Exception as e:
            print(f"Warning: Could not add two-column content: {e}")
    
    return slide


def set_theme_fonts(prs, heading_font_name, body_font_name):
    """Set the theme-level default fonts in the Presentation.
    This ensures even inherited/unformatted text uses the chosen fonts
    instead of the default Calibri/Calibri Light."""
    try:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        # Access the theme blob through the slide master's relationships
        for slide_master in prs.slide_masters:
            for rel in slide_master.part.rels.values():
                if 'theme' in rel.reltype:
                    theme_part = rel.target_part
                    # Parse the theme XML from the blob
                    theme_xml = etree.fromstring(theme_part.blob)
                    font_scheme = theme_xml.find('.//{%s}fontScheme' % ns)
                    if font_scheme is not None:
                        major = font_scheme.find('{%s}majorFont' % ns)
                        minor = font_scheme.find('{%s}minorFont' % ns)
                        if major is not None:
                            latin = major.find('{%s}latin' % ns)
                            if latin is not None:
                                latin.set('typeface', heading_font_name)
                        if minor is not None:
                            latin = minor.find('{%s}latin' % ns)
                            if latin is not None:
                                latin.set('typeface', body_font_name)
                    # Write modified XML back to the theme part blob
                    theme_part._blob = etree.tostring(theme_xml, xml_declaration=True, encoding='UTF-8', standalone=True)
                    break
    except Exception as e:
        print(f"Warning: Could not set theme fonts: {e}")


# ---------------------------------------------------------------------------
# Font directories where installed fonts live on Windows
# ---------------------------------------------------------------------------
_SYSTEM_FONT_DIR = os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts")
_USER_FONT_DIR = os.path.join(os.environ.get("LOCALAPPDATA", ""), "Microsoft", "Windows", "Fonts")
_THEME_UPLOADS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'theme_uploads')


def _find_font_files(font_name, font_url=None):
    """Find .ttf/.otf files on disk for a given font family name.
    Returns a list of file paths (may be empty).
    
    If font_url is provided (from an admin-uploaded font), the corresponding
    file in theme_uploads/ is returned directly.
    """
    results = []
    if not font_name:
        return results

    # 1) If a font_url is given (admin-uploaded), resolve it first
    if font_url:
        import urllib.parse
        decoded = urllib.parse.unquote(font_url)
        filename = decoded.split('/')[-1]
        fpath = os.path.join(_THEME_UPLOADS_DIR, filename)
        if os.path.isfile(fpath):
            results.append(fpath)
            return results
        # Also try the raw (non-decoded) filename
        raw_filename = font_url.split('/')[-1]
        raw_path = os.path.join(_THEME_UPLOADS_DIR, raw_filename)
        if os.path.isfile(raw_path):
            results.append(raw_path)
            return results

    # 2) Search system / user font directories
    # Google Fonts naming: "Open Sans" -> files named OpenSans-*.ttf
    base = font_name.replace(" ", "")
    for d in (_SYSTEM_FONT_DIR, _USER_FONT_DIR):
        if not os.path.isdir(d):
            continue
        for f in os.listdir(d):
            if f.lower().endswith((".ttf", ".otf")) and base.lower() in f.lower().replace("-", "").replace("_", ""):
                results.append(os.path.join(d, f))

    # 3) Also search theme_uploads/ for admin-uploaded font files
    if not results and os.path.isdir(_THEME_UPLOADS_DIR):
        for f in os.listdir(_THEME_UPLOADS_DIR):
            if f.lower().endswith((".ttf", ".otf")) and base.lower() in f.lower().replace("-", "").replace("_", ""):
                results.append(os.path.join(_THEME_UPLOADS_DIR, f))

    return results


def _classify_font_file(filepath):
    """Determine if a font file is Regular, Bold, Italic, or BoldItalic.
    Returns one of: 'regular', 'bold', 'italic', 'boldItalic'."""
    fname = os.path.basename(filepath).lower()
    is_bold = 'bold' in fname or 'semibold' in fname
    is_italic = 'italic' in fname
    if is_bold and is_italic:
        return 'boldItalic'
    elif is_bold:
        return 'bold'
    elif is_italic:
        return 'italic'
    # Variable fonts without explicit style markers → regular
    return 'regular'


def embed_fonts_in_pptx(pptx_path, font_names):
    """Embed font files into a saved PPTX using direct ZIP manipulation.

    This post-save approach modifies the PPTX file on disk rather than the
    in-memory python-pptx object, which is more reliable and avoids OOXML
    schema ordering issues.

    Args:
        pptx_path: path to the saved .pptx file
        font_names: iterable of font family names to embed
    """
    try:
        import zipfile, shutil, tempfile
        from lxml import etree

        P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        REL_NS = 'http://schemas.openxmlformats.org/package/2006/relationships'
        CT_NS = 'http://schemas.openxmlformats.org/package/2006/content-types'
        FONT_REL_TYPE = (
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/font'
        )

        # Collect font files grouped by family
        # Also look up font_url from database for admin-uploaded fonts
        fonts_by_family = {}
        for name in set(font_names):
            if not name:
                continue
            # Try to find font_url from database for this font name
            _, db_font_url = _lookup_custom_font_from_db(name)
            if not db_font_url:
                # Also try looking up by font_family
                try:
                    from app import get_db_connection
                    conn = get_db_connection()
                    cursor = conn.cursor(dictionary=True)
                    cursor.execute(
                        "SELECT font_url FROM tbl_fonts WHERE font_family LIKE %s AND is_active = TRUE AND font_url IS NOT NULL AND font_url != '' LIMIT 1",
                        (f"%{name}%",)
                    )
                    row = cursor.fetchone()
                    cursor.close()
                    conn.close()
                    if row:
                        db_font_url = row.get('font_url', '')
                except Exception:
                    pass
            files = _find_font_files(name, font_url=db_font_url)
            if files:
                fonts_by_family[name] = files

        if not fonts_by_family:
            return

        # Work on a temp copy, then replace the original
        tmp_fd, tmp_path = tempfile.mkstemp(suffix='.pptx')
        os.close(tmp_fd)

        try:
            with zipfile.ZipFile(pptx_path, 'r') as zin, \
                 zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:

                # Read existing rels for presentation.xml
                rels_path = 'ppt/_rels/presentation.xml.rels'
                rels_xml = zin.read(rels_path)
                rels_root = etree.fromstring(rels_xml)

                # Find the max existing rId
                max_rid = 0
                for rel in rels_root:
                    rid = rel.get('Id', '')
                    if rid.startswith('rId'):
                        try:
                            max_rid = max(max_rid, int(rid[3:]))
                        except ValueError:
                            pass

                # Read presentation.xml
                pres_xml = zin.read('ppt/presentation.xml')
                pres_root = etree.fromstring(pres_xml)

                # Read [Content_Types].xml
                ct_xml = zin.read('[Content_Types].xml')
                ct_root = etree.fromstring(ct_xml)

                font_idx = 0
                total_embedded = 0
                font_data_parts = {}  # part_path -> font_data bytes

                # Build <p:embeddedFontLst> element
                embedded_lst = etree.Element('{%s}embeddedFontLst' % P)

                for family_name, file_paths in fonts_by_family.items():
                    style_files = {}
                    for fpath in file_paths:
                        style = _classify_font_file(fpath)
                        if style not in style_files:
                            style_files[style] = fpath

                    if 'regular' not in style_files and file_paths:
                        style_files['regular'] = file_paths[0]

                    emb_font = etree.SubElement(embedded_lst, '{%s}embeddedFont' % P)
                    font_elem = etree.SubElement(emb_font, '{%s}font' % P)
                    font_elem.set('typeface', family_name)
                    font_elem.set('charset', '0')

                    for style_key in ('regular', 'bold', 'italic', 'boldItalic'):
                        if style_key not in style_files:
                            continue

                        fpath = style_files[style_key]
                        font_idx += 1
                        part_path = f'ppt/fonts/font{font_idx}.fntdata'
                        rid = f'rId{max_rid + font_idx}'

                        try:
                            with open(fpath, 'rb') as f:
                                font_data_parts[part_path] = f.read()

                            # Add relationship
                            rel_elem = etree.SubElement(rels_root, 'Relationship')
                            rel_elem.set('Id', rid)
                            rel_elem.set('Type', FONT_REL_TYPE)
                            rel_elem.set('Target', f'fonts/font{font_idx}.fntdata')

                            # Add style element under embeddedFont
                            style_elem = etree.SubElement(
                                emb_font, '{%s}%s' % (P, style_key)
                            )
                            style_elem.set('{%s}id' % R, rid)
                            total_embedded += 1
                        except Exception as e:
                            print(f"  [Embed] Could not read {fpath}: {e}")

                if total_embedded == 0:
                    os.remove(tmp_path)
                    return

                # Insert <p:embeddedFontLst> at the correct schema position
                # OOXML order: sldMasterIdLst, notesMasterIdLst, handoutMasterIdLst,
                #   sldIdLst, sldSz, notesSz, smartTags, embeddedFontLst, ...
                BEFORE_EMBED = {
                    'sldMasterIdLst', 'notesMasterIdLst', 'handoutMasterIdLst',
                    'sldIdLst', 'sldSz', 'notesSz', 'smartTags',
                }
                insert_idx = 0
                for i, child in enumerate(pres_root):
                    tag = etree.QName(child.tag).localname
                    if tag in BEFORE_EMBED:
                        insert_idx = i + 1
                pres_root.insert(insert_idx, embedded_lst)

                # Set presentation attributes
                pres_root.set('embedTrueTypeFonts', '1')
                pres_root.set('saveSubsetFonts', '1')

                # Ensure .fntdata extension has a content type in [Content_Types].xml
                has_fntdata = any(
                    el.get('Extension') == 'fntdata'
                    for el in ct_root
                    if el.tag == '{%s}Default' % CT_NS
                )
                if not has_fntdata:
                    default_elem = etree.SubElement(ct_root, '{%s}Default' % CT_NS)
                    default_elem.set('Extension', 'fntdata')
                    default_elem.set('ContentType', 'application/x-fontdata')

                # Copy all existing entries to the new ZIP, replacing modified XMLs
                for item in zin.infolist():
                    if item.filename == 'ppt/presentation.xml':
                        new_pres = etree.tostring(
                            pres_root, xml_declaration=True,
                            encoding='UTF-8', standalone=True
                        )
                        zout.writestr(item.filename, new_pres)
                    elif item.filename == rels_path:
                        new_rels = etree.tostring(
                            rels_root, xml_declaration=True,
                            encoding='UTF-8', standalone=True
                        )
                        zout.writestr(item.filename, new_rels)
                    elif item.filename == '[Content_Types].xml':
                        new_ct = etree.tostring(
                            ct_root, xml_declaration=True,
                            encoding='UTF-8', standalone=True
                        )
                        zout.writestr(item.filename, new_ct)
                    else:
                        # Use filename string (not ZipInfo) to avoid CRC mismatch
                        # when re-compressing with different settings
                        data = zin.read(item.filename)
                        zout.writestr(item.filename, data)

                # Add font data parts
                for part_path, data in font_data_parts.items():
                    zout.writestr(part_path, data)

            # Replace original with modified file
            shutil.move(tmp_path, pptx_path)
            print(f"[PPT] Embedded {total_embedded} font file(s) for "
                  f"{len(fonts_by_family)} family(ies) into PPTX")

        except Exception as e:
            # Clean up temp file on error
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
            raise

    except Exception as e:
        print(f"[PPT] Font embedding skipped: {e}")


def create_ppt(content, output_path, theme="cosmic_dark", layout="title_bullets",
               heading_font="calibri", body_font="calibri", font_size="standard",
               custom_title_color=None, presentation_title="", content_depth="balanced",
               subtitle="", add_thank_you=False,
               custom_title_color_hex="", custom_heading_color_hex="", custom_body_color_hex=""):
    """Create a professional PowerPoint presentation with all specified settings"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get configurations - check static themes first, then dynamic color themes, then dynamic templates
    theme_config = THEMES.get(theme)
    template_font = None  # Store template font if available
    
    if not theme_config:
        # Check if it's a dynamic color theme (dynamic_color_*)
        theme_config = get_dynamic_themes().get(theme)
    if not theme_config and theme.startswith('dynamic_'):
        # Check if it's a template (dynamic_* without color_ prefix)
        theme_config = get_template_config(theme)
        # If template has custom font, use it
        if theme_config and theme_config.get('font_family'):
            template_font = theme_config.get('font_family')
    if not theme_config:
        theme_config = THEMES["cosmic_dark"]
    
    # Debug: print background info for troubleshooting
    print(f"[PPT] Background type: {theme_config.get('background_type')}, Background image: {theme_config.get('background_image')}")
    
    # Use template font if available, otherwise use provided fonts
    # PRIORITY: User's explicit font selection > Template font > Defaults
    # If user changed heading/body font from defaults, respect their choice
    user_chose_custom_font = (heading_font != "calibri" or body_font != "calibri")
    
    if user_chose_custom_font:
        # User explicitly selected fonts in the dropdown — always respect that
        font_config = get_font_name(heading_font)
        font_config_body = get_font_name(body_font)
        font_config = {
            "heading": font_config.get("heading", "Calibri"),
            "body": font_config_body.get("body", "Calibri"),
            "_heading_font_url": font_config.get("_font_url", ""),
            "_body_font_url": font_config_body.get("_font_url", ""),
        }
    elif template_font:
        # User didn't change fonts, use the template's actual font name
        clean_font = resolve_system_font(template_font)
        font_config = {
            "heading": clean_font,
            "body": clean_font,
            "_heading_font_url": theme_config.get("font_url", ""),
            "_body_font_url": theme_config.get("font_url", ""),
        }
    else:
        font_config = get_font_name(heading_font)
        font_config_body = get_font_name(body_font)
        font_config = {
            "heading": font_config.get("heading", "Calibri"),
            "body": font_config_body.get("body", "Calibri"),
            "_heading_font_url": font_config.get("_font_url", ""),
            "_body_font_url": font_config_body.get("_font_url", ""),
        }
    size_config = FONT_SIZES.get(font_size, FONT_SIZES["standard"])
    
    # Set theme-level default fonts so inherited text also uses the chosen fonts
    set_theme_fonts(prs, font_config['heading'], font_config['body'])
    
    print(f"[PPT] Theme: {theme}, Heading font: {font_config['heading']}, Body font: {font_config['body']}, Size: {font_size}")
    
    # Get custom title color if specified (for title slide)
    title_color = None
    if custom_title_color_hex and custom_title_color_hex.startswith('#'):
        try:
            title_color = hex_to_rgb(custom_title_color_hex)
        except Exception:
            title_color = None
    if not title_color and custom_title_color and custom_title_color in FONT_COLORS:
        title_color = FONT_COLORS[custom_title_color]
    elif not title_color and custom_title_color and custom_title_color.startswith('dynamic_color_'):
        title_color = resolve_dynamic_color(custom_title_color)
    
    if not title_color:
        title_color = theme_config["title_color"]
    
    # Get custom heading color (for content slide headings)
    heading_color = None
    if custom_heading_color_hex and custom_heading_color_hex.startswith('#'):
        try:
            heading_color = hex_to_rgb(custom_heading_color_hex)
        except Exception:
            heading_color = None
    if not heading_color:
        heading_color = theme_config["title_color"]
    
    # Get custom body color (for content slide body text)
    body_color = None
    if custom_body_color_hex and custom_body_color_hex.startswith('#'):
        try:
            body_color = hex_to_rgb(custom_body_color_hex)
        except Exception:
            body_color = None
    if not body_color:
        body_color = theme_config["body_color"]
    
    print(f"[PPT] Title color: {title_color}, Heading color: {heading_color}, Body color: {body_color}")
    
    # =============================================
    # MODE 1: MANUAL CONTENT (LIST OF STRINGS)
    # =============================================
    if isinstance(content, list):
        # Create a working copy of theme_config with custom color overrides
        working_theme = dict(theme_config)
        working_theme["title_color"] = title_color
        working_theme["body_color"] = body_color
        
        # Create a separate theme for content slides where heading_color may differ
        content_theme = dict(working_theme)
        content_theme["title_color"] = heading_color  # Content slide headings use heading_color
        
        # Add title slide first (uses title_color for the main title)
        create_title_slide(
            prs,
            presentation_title if presentation_title else "Presentation",
            subtitle,
            working_theme, font_config, size_config
        )
        
        # Add content slides (uses heading_color for slide titles, body_color for bullets)
        for slide_text in content:
            lines = [l.strip() for l in slide_text.split("\n") if l.strip()]
            if not lines:
                continue
            
            title = lines[0][:80]
            bullets = [l for l in lines[1:] if l]
            
            # Limit bullets to prevent overflow - max 6 per slide
            max_bullets_per_slide = 6
            if len(bullets) > max_bullets_per_slide:
                bullets = bullets[:max_bullets_per_slide]
            
            # NEVER use two-column layout - always stack bullets vertically
            layout_type = "bullets"
            
            create_content_slide(
                prs, title, bullets,
                content_theme, font_config, size_config, layout_type
            )
        
        # Add thank you slide if requested
        if add_thank_you:
            create_thankyou_slide(prs, working_theme, font_config, size_config)
        
        prs.save(output_path)
        # Embed font files into the saved PPTX so it works on any machine
        embed_fonts_in_pptx(output_path, [font_config['heading'], font_config['body']])
        return
    
    # =============================================
    # MODE 2: AI GENERATED CONTENT
    # =============================================
    
    import re
    
    # Parse AI content - look for slide separators
    slides_data = []
    
    # Normalize the content - remove extra whitespace and markdown
    content = content.replace("**", "").replace("##", "").replace("# ", "")
    
    # Try SLIDE_BREAK first (our standard format)
    if "SLIDE_BREAK" in content:
        slide_blocks = content.split("SLIDE_BREAK")
    elif "===" in content:
        slide_blocks = content.split("===")
    elif "---" in content:
        slide_blocks = content.split("---")
    else:
        # Split by double newlines if no clear separator
        slide_blocks = re.split(r'\n\s*\n\s*\n', content)
    
    for block in slide_blocks:
        block = block.strip()
        if not block or len(block) < 10:
            continue
        
        lines = [l.strip() for l in block.split("\n") if l.strip()]
        if len(lines) < 1:
            continue
        
        # First non-empty line is title
        title = lines[0]
        
        # Clean up title - remove common prefixes
        title = re.sub(r'^(Slide\s*\d*:?\s*|Title:?\s*|\d+\.\s*|\[\s*|\]\s*)', '', title, flags=re.IGNORECASE)
        title = title.strip('[]').strip()
        
        # Skip if title looks like a bullet point
        if title.startswith(("•", "-", "–", "—", "*")):
            # Maybe this block has no clear title, use first bullet as title
            title = re.sub(r'^[•\-\–\—\*]\s*', '', title).strip()
        
        # Collect bullet points
        bullets = []
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
            
            # Check if line starts with bullet marker
            if line.startswith(("•", "-", "–", "—", "*")) or re.match(r'^\d+\.', line):
                # Clean the bullet
                clean_line = re.sub(r'^[•\-\–\—\*\d\.]+\s*', '', line).strip()
                if clean_line and len(clean_line) > 3:
                    bullets.append(clean_line)
            elif len(line) > 10 and not line.startswith(("Note:", "Example:", "[", "(")):
                # Long line without bullet marker - treat as content
                bullets.append(line)
        
        # Add slide if we have a valid title
        if title and len(title) > 2:
            # Limit title length
            if len(title) > 80:
                title = title[:77] + "..."
            
            # If no bullets, skip this slide (it's probably noise)
            if not bullets and len(lines) > 1:
                # Try to extract bullets from remaining lines more aggressively
                for line in lines[1:]:
                    clean = line.strip().lstrip("•-–—*0123456789. ")
                    if clean and len(clean) > 5:
                        bullets.append(clean)
            
            if bullets:
                slides_data.append({"title": title, "bullets": bullets})
    
    # Ensure we have at least one slide
    if not slides_data:
        lines = content.split("\n")
        current_title = ""
        current_bullets = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check if it's a title line (short, no bullet)
            if len(line) < 60 and not line.startswith(("•", "-", "–", "—")):
                if current_title and current_bullets:
                    slides_data.append({"title": current_title, "bullets": current_bullets[:]})
                current_title = line
                current_bullets = []
            else:
                # It's a bullet point
                clean_line = line.lstrip("•-–—0123456789. ").strip()
                if clean_line:
                    current_bullets.append(clean_line)
        
        # Add last slide
        if current_title:
            if not current_bullets:
                current_bullets = ["Content will appear here"]
            slides_data.append({"title": current_title, "bullets": current_bullets})
    
    # Ensure at least one Slide
    if not slides_data:
        slides_data = [{"title": "Presentation", "bullets": ["Content will appear here"]}]
    
    # Create title slide
    # Apply custom color overrides for AI mode too
    working_theme_ai = dict(theme_config)
    working_theme_ai["title_color"] = title_color
    working_theme_ai["body_color"] = body_color
    
    content_theme_ai = dict(working_theme_ai)
    content_theme_ai["title_color"] = heading_color
    
    create_title_slide(
        prs,
        presentation_title if presentation_title else slides_data[0]["title"] if slides_data else "Presentation",
        subtitle,
        working_theme_ai, font_config, size_config
    )
    
    # Create content slides
    for slide_data in slides_data[1:] if len(slides_data) > 1 else slides_data:
        title = slide_data.get("title", "Slide")
        bullets = slide_data.get("bullets", ["Content"])
        
        # Limit bullets to max 6 — content fits the slide
        if len(bullets) > 6:
            bullets = bullets[:6]
        
        # ALWAYS use single column vertical layout
        layout_type = "bullets"
        
        create_content_slide(
            prs, title, bullets,
            content_theme_ai, font_config, size_config, layout_type
        )
    
    # Add thank you slide if requested
    if add_thank_you:
        create_thankyou_slide(prs, working_theme_ai, font_config, size_config)
    
    prs.save(output_path)
    # Embed font files into the saved PPTX so it works on any machine
    embed_fonts_in_pptx(output_path, [font_config['heading'], font_config['body']])
